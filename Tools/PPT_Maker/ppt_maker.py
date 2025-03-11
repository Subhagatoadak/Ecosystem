import sys
import os
# Compute the absolute path to the root directory
root_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
if root_path not in sys.path:
    sys.path.insert(0, root_path)

from llm_service.llm_generator import generate_llm_response
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
import io

# For charts
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Define slide layout options (these indices may vary with your template)
layout_options = {
    "Title Slide (0)": 0,
    "Title and Content (1)": 1,
    "Section Header (2)": 2,
    "Two Content (3)": 3,
    "Comparison (4)": 4,
    "Title Only (5)": 5,
    "Blank (6)": 6,
    "Content with Caption (7)": 7,
    "Picture with Caption (8)": 8,
    "Title and Vertical Text (9)": 9,
    "Vertical Title and Text (10)": 10
}

# Define chart type options.
chart_type_options = {
    "Column Clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "Bar Clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "Line": XL_CHART_TYPE.LINE,
    "Pie": XL_CHART_TYPE.PIE,
    "Scatter": XL_CHART_TYPE.XY_SCATTER
}

def create_presentation(presentation_title, description, author,
                        title_bg_bytes, common_content_bg_bytes,
                        sections_data):
    prs = Presentation()
    
    # (Optional) Print available layouts for debugging.
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}: {layout.name}")
    
    # ----------------------------
    # Create Main Title Slide
    # ----------------------------
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    if title_bg_bytes:
        bg_stream = io.BytesIO(title_bg_bytes)
        bg = slide.shapes.add_picture(bg_stream, 0, 0,
                                      width=prs.slide_width,
                                      height=prs.slide_height)
        # Move image behind other shapes.
        bg._element.getparent().remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)
    
    slide.shapes.title.text = presentation_title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"{description}\n\nAuthor: {author}"
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2),
                                         prs.slide_width - Inches(2),
                                         Inches(1))
        txBox.text = f"{description}\n\nAuthor: {author}"
    
    # ----------------------------
    # Process Each Section
    # ----------------------------
    for section in sections_data:
        section_title = section["section_title"]
        section_header_bg = section.get("section_header_bg", None)
        slides = section["slides"]
        
        # Create section header slide.
        try:
            section_layout = prs.slide_layouts[2]
        except IndexError:
            section_layout = prs.slide_layouts[0]
        sec_slide = prs.slides.add_slide(section_layout)
        if sec_slide.shapes.title:
            sec_slide.shapes.title.text = section_title
        else:
            txBox = sec_slide.shapes.add_textbox(Inches(1), Inches(1),
                                                 prs.slide_width - Inches(2),
                                                 Inches(1))
            txBox.text = section_title
        
        # Add background for section header if provided.
        if section_header_bg:
            bg_stream = io.BytesIO(section_header_bg)
            bg = sec_slide.shapes.add_picture(bg_stream, 0, 0,
                                              width=prs.slide_width,
                                              height=prs.slide_height)
            bg._element.getparent().remove(bg._element)
            sec_slide.shapes._spTree.insert(2, bg._element)
        
        # Create slides for this section.
        for idx, slide_data in enumerate(slides):
            layout_index = slide_data.get("layout", 6)
            content = slide_data.get("content", "")
            image_bytes = slide_data.get("image", None)
            image_type = slide_data.get("image_type", None)  # "background" or "foreground"
            chart_type = slide_data.get("chart_type", None)
            font_size = slide_data.get("font_size", 24)
            font_type = slide_data.get("font_type", "Calibri")
            
            try:
                slide_layout = prs.slide_layouts[layout_index]
            except IndexError:
                slide_layout = prs.slide_layouts[6]
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Set a default title for the slide if available.
            if new_slide.shapes.title:
                new_slide.shapes.title.text = f"{section_title} - Slide {idx+1}"
            
            # Add text content and apply font settings.
            if content:
                if len(new_slide.placeholders) > 1:
                    placeholder = new_slide.placeholders[1]
                    placeholder.text = content
                    text_frame = placeholder.text_frame
                else:
                    txBox = new_slide.shapes.add_textbox(Inches(1), Inches(2),
                                                         prs.slide_width - Inches(2),
                                                         Inches(2))
                    txBox.text = content
                    text_frame = txBox.text_frame
                # Apply font formatting.
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                        run.font.name = font_type
            
            # Determine which image to use:
            # Use slide-specific image if provided; otherwise, use common background if available.
            use_image_bytes = image_bytes if image_bytes else common_content_bg_bytes
            
            if use_image_bytes:
                if image_bytes and image_type == "foreground":
                    # For a foreground image, add a smaller picture.
                    bg_stream = io.BytesIO(use_image_bytes)
                    # Position the image at bottom-right with a width of 3 inches.
                    x = prs.slide_width - Inches(3) - Inches(0.5)  # 0.5 inch margin
                    y = prs.slide_height - Inches(3) - Inches(0.5)
                    new_slide.shapes.add_picture(bg_stream, x, y, width=Inches(3))
                else:
                    # For background images (or if using the common background), add full-slide image.
                    bg_stream = io.BytesIO(use_image_bytes)
                    pic = new_slide.shapes.add_picture(bg_stream, 0, 0,
                                                       width=prs.slide_width,
                                                       height=prs.slide_height)
                    pic._element.getparent().remove(pic._element)
                    new_slide.shapes._spTree.insert(2, pic._element)
            
            # Add a chart if requested.
            if chart_type:
                chart_const = chart_type_options.get(chart_type, None)
                if chart_const:
                    chart_data = CategoryChartData()
                    chart_data.categories = ['A', 'B', 'C']
                    chart_data.add_series('Series 1', (10, 20, 30))
                    # Position the chart (customize as needed)
                    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
                    new_slide.shapes.add_chart(chart_const, x, y, cx, cy, chart_data)
    
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

def main():
    st.title("Advanced PPT Creator App")
    st.write("Configure your presentation details below.")
    
    # --- Basic Presentation Details ---
    presentation_title = st.text_input("Presentation Title", "My Presentation")
    description = st.text_area("Description", "This is a description for the presentation.")
    author = st.text_input("Author", "John Doe")
    
    # --- Title Slide Background ---
    add_title_bg = st.checkbox("Add a background image for the title slide?")
    title_bg_bytes = None
    if add_title_bg:
        title_bg_file = st.file_uploader("Upload title slide background", type=["png", "jpg", "jpeg"], key="title_bg")
        if title_bg_file is not None:
            st.image(title_bg_file, caption="Title Slide Background", use_column_width=True)
            title_bg_bytes = title_bg_file.getvalue()
    
    # --- Common Background for Content Slides ---
    add_common_bg = st.checkbox("Add a common background image for all content slides?")
    common_content_bg_bytes = None
    if add_common_bg:
        common_bg_file = st.file_uploader("Upload common background for slides", type=["png", "jpg", "jpeg"], key="common_bg")
        if common_bg_file is not None:
            st.image(common_bg_file, caption="Common Content Background", use_column_width=True)
            common_content_bg_bytes = common_bg_file.getvalue()
    
    # --- Section Header Background Images ---
    add_section_bg = st.checkbox("Add background images for section header slides?")
    
    # --- Sections & Slides ---
    use_sections = st.checkbox("Do you want to create sections?")
    sections_data = []
    if use_sections:
        num_sections = st.number_input("Number of Sections", min_value=1, step=1, value=1)
        for s in range(int(num_sections)):
            with st.expander(f"Section {s+1} Details", expanded=True):
                section_title = st.text_input(f"Section {s+1} Title", f"Section {s+1}", key=f"section_title_{s}")
                section_header_bg = None
                if add_section_bg:
                    sec_bg_file = st.file_uploader(f"Upload background for Section {s+1} header", type=["png", "jpg", "jpeg"], key=f"sec_bg_{s}")
                    if sec_bg_file is not None:
                        st.image(sec_bg_file, caption=f"Section {s+1} Header Background", use_column_width=True)
                        section_header_bg = sec_bg_file.getvalue()
                num_slides = st.number_input(f"Number of slides in Section {s+1}", min_value=0, step=1, value=1, key=f"num_slides_{s}")
                add_content = st.checkbox(f"Add content, images, or charts to slides in Section {s+1}?", key=f"add_content_{s}")
                slides = []
                if num_slides > 0:
                    slide_tabs = st.tabs([f"Slide {i+1}" for i in range(int(num_slides))])
                    for i, tab in enumerate(slide_tabs):
                        with tab:
                            layout_choice = st.selectbox(
                                f"Select layout for Slide {i+1}",
                                list(layout_options.keys()),
                                key=f"layout_{s}_{i}"
                            )
                            content = ""
                            image_bytes = None
                            image_type = None
                            chart_type = None
                            use_ai = False
                            ai_prompt = ""
                            font_size = 24  # default
                            font_type = "Calibri"  # default
                            if add_content:
                                content = st.text_area(f"Content for Slide {i+1}", key=f"content_{s}_{i}")
                                # Checkbox for using AI to rewrite content.
                                use_ai = st.checkbox(f"Use AI to rewrite content for Slide {i+1}?", key=f"use_ai_{s}_{i}")
                                if use_ai:
                                    ai_prompt = st.text_area("Enter AI prompt for rewriting:", key=f"ai_prompt_{s}_{i}")
                                # Options for font size and type.
                                font_size = st.number_input("Font Size", min_value=8, max_value=72, value=24, key=f"font_size_{s}_{i}")
                                font_type = st.selectbox("Font Type", options=["Calibri", "Arial", "Times New Roman", "Verdana", "Comic Sans MS"], key=f"font_type_{s}_{i}")
                                add_slide_image = st.checkbox(f"Add an image for Slide {i+1}?", key=f"add_image_{s}_{i}")
                                if add_slide_image:
                                    image_type = st.radio(f"Image type for Slide {i+1}", options=["background", "foreground"], key=f"img_type_{s}_{i}")
                                    slide_image_file = st.file_uploader(f"Upload image for Slide {i+1}", type=["png", "jpg", "jpeg"], key=f"slide_image_{s}_{i}")
                                    if slide_image_file is not None:
                                        st.image(slide_image_file, caption=f"Slide {i+1} Image", use_column_width=True)
                                        image_bytes = slide_image_file.getvalue()
                                add_chart = st.checkbox(f"Add a chart to Slide {i+1}?", key=f"add_chart_{s}_{i}")
                                if add_chart:
                                    chart_type = st.selectbox(f"Select chart type for Slide {i+1}",
                                                              list(chart_type_options.keys()),
                                                              key=f"chart_{s}_{i}")
                            slides.append({
                                "layout": layout_options[layout_choice],
                                "content": content,
                                "image": image_bytes,
                                "image_type": image_type,
                                "chart_type": chart_type,
                                "use_ai": use_ai,
                                "ai_prompt": ai_prompt,
                                "font_size": font_size,
                                "font_type": font_type
                            })
                sections_data.append({
                    "section_title": section_title,
                    "section_header_bg": section_header_bg,
                    "slides": slides
                })
    else:
        st.info("No sections selected. A default section with one slide will be added.")
        sections_data.append({
            "section_title": "Default Section",
            "section_header_bg": None,
            "slides": [{
                "layout": layout_options["Title and Content (1)"],
                "content": "",
                "image": None,
                "image_type": None,
                "chart_type": None,
                "use_ai": False,
                "ai_prompt": "",
                "font_size": 24,
                "font_type": "Calibri"
            }]
        })
    
    if st.button("Generate PPT"):
        # Iterate over sections and slides. If AI rewriting is requested, update the content.
        for section in sections_data:
            for slide_data in section["slides"]:
                if slide_data.get("use_ai", False):
                    original_content = slide_data.get("content", "")
                    ai_prompt = slide_data.get("ai_prompt", "")
                    if original_content and ai_prompt:
                        # Call the LLM service to rewrite the content.
                        slide_data["content"] = generate_llm_response(
                            "Context:\n" + original_content + "\n\n" + "Instructions:\n" + ai_prompt,
                            provider="openai",
                            model="gpt-4o",
                            temperature=0.7
                        )
        ppt_file = create_presentation(presentation_title, description, author,
                                       title_bg_bytes, common_content_bg_bytes,
                                       sections_data)
        st.success("Presentation generated successfully!")
        st.download_button(
            label="Download PPT",
            data=ppt_file,
            file_name="advanced_generated_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if __name__ == "__main__":
    main()
